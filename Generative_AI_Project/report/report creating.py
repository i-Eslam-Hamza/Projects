from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a PowerPoint presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Global Happiness Report 2016 Dashboard"
subtitle.text = "Analysis and Insights"

# Define a consistent style for headings and text
def add_heading(slide, text):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True

def add_text(slide, text):
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)

def add_plot_placeholder(slide, placeholder_text):
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(8), Inches(1))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = placeholder_text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)

# Slide 2: Project Purpose Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_heading(slide, "Project Purpose")
add_text(slide, "The purpose of this analysis was to identify and understand key factors contributing to happiness worldwide, using data from the 2016 World Happiness Report. The goal is to offer insights into how economic, health, and social indicators impact national well-being.")

# Slide 3: Top 10 Countries - GDP and Life Expectancy
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_heading(slide, "Top 10 Countries by GDP per Capita and Life Expectancy")
add_text(slide, "The bar charts display the top 10 happiest countries based on GDP per capita and healthy life expectancy, revealing the role these indicators play in happiness.")
add_plot_placeholder(slide, "Insert Bar Charts: Top 10 Countries")

# Slide 4: Happiness and GDP per Capita by Region
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_heading(slide, "Happiness and GDP per Capita by Region")
add_text(slide, "This scatter plot highlights the relationship between GDP per capita and Happiness Score across regions, showing how economic resources contribute to happiness in different parts of the world.")
add_plot_placeholder(slide, "Insert Scatter Plot: GDP vs. Happiness Score")

# Slide 5: Happiness Score Distribution by Region
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_heading(slide, "Happiness Score Distribution by Region")
add_text(slide, "The pie chart shows the percentage contribution of each region to the Happiness Score, providing a visual assessment of global happiness distribution.")
add_plot_placeholder(slide, "Insert Pie Chart: Happiness Score by Region")

# Slide 6: Global Map - GDP per Capita and Life Expectancy by Country
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_heading(slide, "Global Map: GDP per Capita and Life Expectancy by Country")
add_text(slide, "The map provides an interactive view of GDP per capita with Healthy Life Expectancy as a tooltip, highlighting countries with higher economic and health standards.")
add_plot_placeholder(slide, "Insert World Map")

# Slide 7: Key Takeaways
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_heading(slide, "Key Takeaways")
add_text(slide, "1. Economic and Health Indicators: High GDP and life expectancy correlate with higher happiness scores.\n2. Regional Disparities: Wealthier regions show higher happiness, while regions with limited resources report lower levels.\n3. Governance and Trust: Trust in government plays a role in national happiness.")

# Slide 8: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_heading(slide, "Conclusion")
add_text(slide, "This dashboard provides insights into global happiness, emphasizing the importance of economic stability, health, and governance. The analysis offers valuable information for policymakers aiming to address disparities and improve global well-being.")

# Save the presentation
prs.save("Happiness_Report_Presentation.pptx")
